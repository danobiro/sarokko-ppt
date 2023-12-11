from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab import rl_config
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from io import BytesIO

import os
import copy
import requests

from time import sleep

"""TODO list:
* Handle various forms of the input of Bible vers locations
* After that, do pretty formatting
* Make a simple GUI interface
* Handle the API request error
* Find out the ratio between normal and superscript numbers
"""

# Decorator to fix bg and text colors
def copy_slide(copyFromPres, slideIndex,  pasteIntoPres):
    slide = Slide(SlideCopyFromPasteInto(copyFromPres, slideIndex,  pasteIntoPres))
    # Give them black background
    slide.fill_bg_solid(0,0,0)
    # Give them white text color?

# Code from https://stackoverflow.com/a/73954830, made by Josip Pardon
# Improved with jimmiesrustled's answer
def SlideCopyFromPasteInto(copyFromPres, slideIndex,  pasteIntoPres):

    # specify the slide you want to copy the contents from
    slide_to_copy = copyFromPres.slides[slideIndex]

    # Define the layout you want to use from your generated pptx

    slide_layout = pasteIntoPres.slide_layouts.get_by_name("Blank") # names of layouts can be found here under step 3: https://www.geeksforgeeks.org/how-to-change-slide-layout-in-ms-powerpoint/
    # it is important for slide_layout to be blank since you dont want these "Write your title here" or something like that textboxes
    # alternative: slide_layout = pasteIntoPres.slide_layouts[copyFromPres.slide_layouts.index(slide_to_copy.slide_layout)]
    
    # create now slide, to copy contents to 
    new_slide = pasteIntoPres.slides.add_slide(slide_layout)

    # create images dict
    imgDict = {}

    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated
    num = 0
    for shp in slide_to_copy.shapes:
        #if 'Picture' in shp.name:
        if hasattr(shp, 'image'):
            # save image
            #shp.name
            pic_name = str(num)
            with open(pic_name+'.jpg', 'wb') as f:
                f.write(shp.image.blob)

            # add image to dict
            imgDict[pic_name+'.jpg'] = [shp.left, shp.top, shp.width, shp.height]
            num += 1
        else:
            # create copy of elem
            el = shp.element
            newel = copy.deepcopy(el)

            # add elem to shape tree
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    # things added first will be covered by things added last => since I want pictures to be in foreground, I will add them after others elements
    # you can change this if you want
    # add pictures
    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

    return new_slide # this returns slide so you can instantly work with it when it is pasted in presentation

# Code from https://stackoverflow.com/a/61922454
def set_subscript(font):
    font._element.set('baseline', '-25000')

def set_superscript(font):
    font._element.set('baseline', '30000')

def set_strikethrough(font):
    font._element.set('strike','sngStrike')


def get_text_width(text, font_size, font_name='Calibri'):
    # Create a PDF document
    pdf_buffer = BytesIO()
    pdf_canvas = canvas.Canvas(pdf_buffer, pagesize=letter)
    pdf_canvas.setFont(font_name, font_size)

    # Draw the text on the PDF
    pdf_text_width = pdf_canvas.stringWidth(text, font_name, font_size)

    # Close the PDF
    pdf_canvas.save()

    return Inches(pdf_text_width / 72.0)  # Convert from points to inches (1 inch = 72 points)

def get_next_stop(vers_text,font_size,font_type,max_width):
    words_len = 0
    words_list = vers_text.split(" ")
    #curr_words = words_list[0]
    curr_words = ""
    #stop_ind = len(curr_words)
    # We will never quit after the first word, so this is fine
    stop_ind = len(words_list[0])

    delimiter_count = 0

    c = 0

    #delimiter_width = get_text_width('@',font_size,font_type)
    
    # superscript ratio - adjust it for better results
    ss_ratio = 0.5
    
    for word in words_list:
        # Remove delimitters for the measurement
        curr_del_c = word.count('@')
        delimiter_count += curr_del_c

        # make corrections for the numbers that will be made superscripts
        correction = 0
        new_word = word
        for j in range(curr_del_c):
            num = new_word.split('@')[1][0]
            correction += (1-ss_ratio) * get_text_width(num,font_size,font_type)
            new_word = new_word.replace('@',"",1)

        # Don't append space before the first word
        if c > 0:
            curr_words += " " + new_word
        else:
            curr_words += new_word
            c += 1
        words_len = get_text_width(curr_words,font_size,font_type) - correction
        #print(words_len,max_width,stop_ind,curr_words)
        #print(words_len/914400)
        
        if words_len > max_width:
            return stop_ind
        else:
            # Add them back as they are still part of the full string
            stop_ind = len(curr_words) + delimiter_count
        
    return stop_ind

def get_scaling_factor(width):
    # Note: in Libre Office you need a different number!
    # This was calculated for MS PowerPoint
    # This info can be used to calculate the scaling ratio
    txt = 'a'*44

    return width/get_text_width(txt,VersContentTextBox.font_size,'Calibri')

# The textbox provided by pptx has weird behavior, but according to
# the documentation it should work similarly to this
class TextBox:
    #def __init__(self,slide,left,top,width,height,text=None,font='Calibri',font_size=28):
    def __init__(self,slide,left,top,width,height,text=None):
        tbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        self.tbox = tbox
        # Make the rectangle transparent
        self.tbox.fill.background()
        self.tbox.line.fill.background()

        self.tf = self.tbox.text_frame
        self.p = self.tf.paragraphs[0]
        self.p.alignment = PP_ALIGN.CENTER
        self.run = self.p.add_run()

        if text is not None:
            #self.set_text(text,font,font_size)
            self.set_text(text)

    #def set_text(self,text,font='Calibri',font_size=28):
    def set_text(self,text):
        #self.text = text
        run = self.run
        run.text = text
        #_font = run.font
        #_font.name = font
        #_font.size = Pt(font_size)
        #_font.bold = False
        
    def set_alignment(self,align):
        p = self.p
        if align == 'left':
            p.alignment = PP_ALIGN.LEFT
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.CENTER

    def make_bold(self):
        run = self.run
        _font = run.font
        _font.bold = True

    def set_font(self,font):
        run = self.run
        _font = run.font
        _font.name = font
        _font.bold = False

    def set_font_size(self,font_size):
        run = self.run
        _font = run.font
        _font.size = Pt(font_size)

    def set_font_color(self,r,g,b):
        run = self.run
        _font = run.font
        _font.color.rgb = RGBColor(r,g,b)

class VersTextBox(TextBox):
    def __init__(self,slide,left,top,width,height):
        super().__init__(slide,left,top,width,height,text=None)
        #self.text_box = TextBox(slide,left,top,width,height,vers_place)
        self.set_font('Calibri')
        # White
        self.set_font_color(255,255,255)
        

class VersPlaceTextBox(VersTextBox):
    # Static variables
    font_size = 28
    
    def __init__(self,slide,left,top,width,height,vers_place):
        super().__init__(slide,left,top,width,height)
        # vers_place may need to be adjusted here if abrevations are allowed
        self.set_font_size(VersPlaceTextBox.font_size)
        self.make_bold()
        self.set_alignment('right')
        self.set_text(vers_place)

    def set_text(self,vers_place):
        #TODO: there can be various user inputs, they have to be handled somewhere. This is just wishful thinking that it will be an input like this
        parts = vers_place.split(" ")
        book = parts[0]
        text = BibleBooks.BOOKS_DICT[book] + " "
        for i in range(1,len(parts)):
            # Convention of how vers places are denoted
            p = parts[i].replace(',',':')
            text += p + " "
        
        super().set_text(text)
        
class VersContentTextBox(VersTextBox):
    font_size = 32

    def __init__(self,slide,left,top,width,height,vers_cont,prs):
        super().__init__(slide,left,top,width,height)
        self.set_run_properties()

        self.max_lines = 4

        self.set_multiline_text(vers_cont,prs)

    def set_run_properties(self):
        self.set_font_size(VersContentTextBox.font_size)
        self.set_alignment('left')
        # Write text content
        # Safety measure, shouldn't be needed if scaling is right
        self.tf.word_wrap = True
        # Mulitple lines, anchor text to the top of the box
        self.tf.vertical_anchor = MSO_ANCHOR.TOP
        # You can't have more than 4 lines on the slide
        #TODO: change this to 4 and a quarter

    def add_text(self,text):
        if '@' not in text:
            self.run.text += text
        else:
            parts = text.split('@')

            for i,part in enumerate(parts):
                # The first part doesn't contain a number
                if i == 0:
                    self.run.text += part
                else:
                    # Create a new run for superscript text
                    self.run = self.p.add_run()
                    self.set_run_properties()
                    font = self.run.font
                    set_superscript(font)
                    # first element is the number
                    self.run.text += part[0]

                    # Return to normal state and add the rest
                    self.run = self.p.add_run()
                    self.set_run_properties()
                    self.run.text += part[1:]


    def add_text_superscript(self,text):
        pass
        
    def set_multiline_text(self,vers_text,prs):
        num_lines = 0
        
        #TODO: this line is dirty
        vers_text_width = get_text_width(vers_text, VersContentTextBox.font_size, 'CalibriBd')

        width = self.tbox.width
        vers_font_size = VersContentTextBox.font_size

        new_slide = False

        while(vers_text_width > width):
            # TODO: this line is dirty
            stop_ind = get_next_stop(vers_text,vers_font_size,'CalibriBd', width / BibleVersSlide.SCALING_FACTOR)
            self.add_text(vers_text[:stop_ind] + '\n')
            #self.run.text += vers_text[:stop_ind] + '\n'

            # Drop spaces
            vers_text = vers_text[stop_ind+1:]
            
            # Add text
            num_lines += 1
            
            # If we end with 10 lines, 1 more will be appended...
            if num_lines >= self.max_lines:
                new_slide = True
            
                # We should keep track of the vers we are at
                # so that we can add that on the next slide

                #TODO: do smart formatting here, such as look back a few words whether we have some punctuation
                # and then start next slide from there
                
                BibleVersSlide(prs,vers_place,vers_text)
                
                # To debug scaling
                #print()
                
                break
                
            vers_text_width = get_text_width(vers_text, vers_font_size, 'CalibriBd')

        if not new_slide:
            #TODO: Make superscript!
            self.run.text += vers_text

        # To debug scaling:
        #print(self.run.text)

        

# Abstract class, no instance should be created, unless from already existing slides
class Slide:
    #@abstractmethod
    def __init__(self,prs):
        self.slide = None
        pass

    def __init__(self,slide):
        self.slide = slide

    def fill_bg_solid(self, r, g, b):
        # Set black background
        if self.slide is None:
            print("Error: Slide is an abstract class, no instance should be created of it")
            quit()
        background = self.slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

class BlankSlide(Slide):
    def __init__(self,prs):
        super().__init__(prs)
        title_only_slide_layout = prs.slide_layouts[6]
        self.slide = prs.slides.add_slide(title_only_slide_layout)

class BibleVersSlide(BlankSlide):
    # This is a constant after the value is determined
    SCALING_FACTOR = 0.98

    def __init__(self,prs,vers_place,vers_cont):
        super().__init__(prs)
        # Global offset for Bible vers slides from the edges of the slide
        self.OFFSET = Cm(0.27)
        OFFSET = self.OFFSET
        # Slide effective width
        self.slide_w = prs.slide_width - OFFSET
        slide_w = self.slide_w
        # Black background
        self.fill_bg_solid(0,0,0)

        # Get vers place textbox
        left = OFFSET
        top = Cm(0.45)
        height = Cm(1.6)
        width = slide_w - left

        # Scaling factor 1.0 seems to work for now, but it may need adjustments
        #BibleVersSlide.SCALING_FACTOR = get_scaling_factor(width)
        
        self.vers_place_box = VersPlaceTextBox(self.slide,left,top,width,height,vers_place)

        # Get vers content textbox

        #top = OFFSET + Pt(VersPlaceTextBox.font_size) + OFFSET
        #height = prs.slide_height - OFFSET - top
        #width = slide_w - left
        left = OFFSET
        top = Cm(1.83)
        height = Cm(11.31)
        width = Cm(24.86)

        self.vers_cont_box = VersContentTextBox(self.slide,left,top,width,height,vers_cont,prs)

def get_bible_books():
    # Get the names into a dict, so we can use it on the ppts
    response = requests.get("https://szentiras.hu/api/books/RÚF")
    # Don't request too frequently from the API!
    sleep(0.5)

    if not response:
        print("Error: couldn't get Bible books")
        quit()

    books = response.json()['books']
    bible_books = {}

    for book in books:
        bible_books[book['abbrev']] = book['name']

    return bible_books

# Could be placed in a different module to make it similar to a Singleton, to prevent it from re-fetching the list
class BibleBooks:
    BOOKS_DICT = {}

    def __init__(self):
        if os.path.isfile("./resources/Bible-books.txt"):
            with open("./resources/Bible-books.txt", 'r', encoding='utf-8') as f:
                for lines in f.readlines():
                    item = lines.split(":")
                    BibleBooks.BOOKS_DICT[item[0].strip()] = item[1].strip()
        else:
            print("Downloading Bible book names")
            BibleBooks.BOOKS_DICT = get_bible_books()
            self.save("Bible-books.txt")

    def save(self,filename):
        # Make sure to include encoding or it won't work!
        with open(f"./resources/{filename}", 'w', encoding='utf-8') as f:
            for key, value in BibleBooks.BOOKS_DICT.items():
                f.write(f"{key} : {value}\n")

def get_vers_content(vers_place):
    """ When vers_place gets here, it should follow the the convention: "Book #,##-##" """
    #TODO: Add support for Book #,##-#,##? (Goes over to the next chapter)
    if vers_place.count(",") > 1:
        #TODO: handle this
        print("ATTENTION: Vers extends a chapter")
    data = vers_place.split(",")
    root = data[0]
    section = data[1]
    section_bounds = section.split("-")
    section_start = int(section_bounds[0].strip())
    section_end = int(section_bounds[1].strip())

    #vers_cont = []
    vers_text = ""

    # inclusive range!
    for i in range(section_start,section_end+1):
        place = root + ',' + str(i)

        success = False
        while(not success):
            response = requests.get(f"https://szentiras.hu/api/ref/{place}/RÚF")

            if not response:
                print("Error: couldn't get Bible verses")
                success = False
                sleep(2)
            else:
                success = True
        
            sleep(0.5)

            # Handle an exception that sometimes arises
            if success == True:
                try:
                    resp_text = response.json()['text']
                except:
                    print("An exception happened. Retrying...")
                    success = False

        #for i in range(len(vers_cont)):
        #    if i != 0:
        #        vers_text += str(i)
        #    vers_text += vers_cont[i]

        #vers_cont.append(resp_text)
        if i != section_start:
            # Use @ as delimitter as that never happens in the text
            vers_text += '@' + str(i)
        vers_text += resp_text

    return vers_text

def create_bible_vers_slides(prs,vers_place):
    #TODO: request vers by vers so we can add vers numbers in the text
    vers_text = get_vers_content(vers_place)
    #response = requests.get(f"https://szentiras.hu/api/ref/{vers_place}/RÚF")
    # Don't request too frequently from the API!
    #sleep(1)
    
    #if not response:
    #    print("Error: couldn't get Bible verses")
    #    quit()

    #vers_cont = response.json()['text']

    #vers_cont = 'Perferendis id voluptatem maxime. Vero debitis dolorem iste blanditiis ut accusamus consectetur omnis. Maiores quasi et rerum voluptate aperiam uti nisi nihil. Quos laborum hic nihil. Nihil perferendis id quia. Minima incidunt molestiae laboriosam ut unde odit quos dolores.…'
    
    # Use this to help find the scaling ratio
    #vers_cont = 'a'*44 + ' ' * 2

    BibleVersSlide(prs,vers_place,vers_text)

def add_song_slides(prs,song_list):
    for song in song_list:
        fname = song + '.pptx'
        song_prs = Presentation(f"./resources/songs/{fname}")

        for i in range(len(song_prs.slides)):
            copy_slide(song_prs, i,  prs)
            

if __name__ == "__main__":
    ##### USER INPUTS #######
    # Songs before Bible vers
    pre_bv_song_list = ['Teremtsd bennem tiszta szívet', 'Teremtsd bennem tiszta szívet']
    vers_place_list = ["Tit 3,3-7"]
    # Songs after Bible vers
    post_bv_song_list = ['Teremtsd bennem tiszta szívet']
    # Post teaching songs (like for communion)
    post_song_list = []
    # Old slide name
    old_slide_name = "pelda.pptx"
    # Index of slide at which ads start in last slide
    ad_start_ind = 95


    # Load Calibri font
    rl_config.TTFSearchPath.append('./resources/calibri-font-family')

    pdfmetrics.registerFont(TTFont('Calibri', 'calibri-regular.ttf'))
    pdfmetrics.registerFont(TTFont('CalibriBd', 'calibri-bold.ttf'))
    pdfmetrics.registerFont(TTFont('CalibriIt', 'calibri-italic.ttf'))
    pdfmetrics.registerFont(TTFont('CalibriBI', 'calibri-bold-italic.ttf'))

    pdfmetrics.registerFontFamily('Calibri',normal='Calibri',bold='CalibriBd',italic='CalibriIt',boldItalic='CalibriBI')

    rl_config.TTFSearchPath.remove('./resources/calibri-font-family')

    # Init
    #BibleBooks().save("Bible-books.txt")
    BibleBooks()

    ####
    # Load old presentation
    old_prs = Presentation(old_slide_name)

    ## Create the presentation
    prs = Presentation()
    # Dimensional convetions
    prs.slide_width = 9144000
    prs.slide_height = 5143500
    # or
    #prs.slide_width = old_prs.slide_width
    #prs.slide_height = old_prs.slide_height

    ## Building the new presentation

    # Create an empty slide for the welcome slide
    BlankSlide(prs).fill_bg_solid(0,0,0)

    # Copy the standard second slide
    copy_slide(old_prs, 1,  prs)

    # Blank black slide
    BlankSlide(prs).fill_bg_solid(0,0,0)

    # Empty black slide for the welcome vers - should it be supported?
    BlankSlide(prs).fill_bg_solid(0,0,0)

    # Blank black slide
    BlankSlide(prs).fill_bg_solid(0,0,0)

    ## Songs before teaching
    # Song before Bible verses
    add_song_slides(prs,pre_bv_song_list)

    # Bible vers slides
    for vers_place in vers_place_list:
        create_bible_vers_slides(prs,vers_place)

    # Blank black slide
    BlankSlide(prs).fill_bg_solid(0,0,0)

    # Add song after Bible verses
    add_song_slides(prs,post_bv_song_list)

    # Blank black slide for teaching slides
    BlankSlide(prs).fill_bg_solid(0,0,0)

    # Post teaching songs
    # if not empty
    if post_song_list:
        add_song_slides(prs,post_song_list)

    # Copy ad slides
    for i in range(ad_start_ind-1, len(old_prs.slides)):
        # And them black bg
        copy_slide(old_prs, i,  prs)

    # Save file
    prs.save('test.pptx')